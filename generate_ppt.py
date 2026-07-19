import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_pitch_deck():
    prs = Presentation()
    
    # 16:9 Widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Layout definitions
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]
    blank_layout = prs.slide_layouts[6]
    
    # Premium Curated Color Palette (Vibrant Neon accents on deep space background)
    SPACE_BLACK = RGBColor(10, 12, 16)      # Deepest background
    PURE_WHITE = RGBColor(240, 246, 252)    # Main text
    CYAN_ELECTRIC = RGBColor(88, 166, 255)  # Actionable headers
    NEON_GREEN = RGBColor(57, 219, 100)     # Positive metrics / highlights
    CHARCOAL_MUTED = RGBColor(139, 148, 158)# Secondary labels
    WARNING_RED = RGBColor(248, 113, 113)   # Issues/cons
    
    def set_premium_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = SPACE_BLACK

    # ----------------------------------------------------
    # Slide 1: Title Slide (Premium Pitch Deck Cover)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(title_layout)
    set_premium_background(slide)
    
    title_box = slide.shapes.title
    title_box.text = "SURFCLAW"
    for paragraph in title_box.text_frame.paragraphs:
        paragraph.font.size = Pt(72)
        paragraph.font.bold = True
        paragraph.font.color.rgb = CYAN_ELECTRIC
        paragraph.font.name = "Segoe UI"
        paragraph.alignment = PP_ALIGN.CENTER
        
    subtitle_box = slide.placeholders[1]
    subtitle_box.text = "The High-Performance AI Agent OS for DePIN Compute Networks\nAccelerated Runtime  |  Hardware-Level MicroVM Isolation  |  Auto-Resiliency"
    for paragraph in subtitle_box.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = CHARCOAL_MUTED
        paragraph.font.name = "Segoe UI"
        paragraph.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # Slide 2: The Problem (Market Pain Points)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(content_layout)
    set_premium_background(slide)
    
    title_box = slide.shapes.title
    title_box.text = "01. The Problem: Critical Pain Points in DePIN Mining"
    title_box.text_frame.paragraphs[0].font.size = Pt(28)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = CYAN_ELECTRIC
    title_box.text_frame.paragraphs[0].font.name = "Segoe UI"
    
    tf = slide.placeholders[1].text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "• Frequent GPU VRAM OOM (Out-of-Memory) Crashes"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = PURE_WHITE
    p.font.name = "Segoe UI"
    p2 = tf.add_paragraph()
    p2.text = "  - Concurrent LLM requests trigger sudden memory overflow, crashing nodes and causing severe downtime."
    p2.font.size = Pt(13)
    p2.font.color.rgb = CHARCOAL_MUTED
    p2.font.name = "Segoe UI"
    
    p3 = tf.add_paragraph()
    p3.text = "• Severe Validation Penalties from Minor Formatting Errors"
    p3.font.size = Pt(15)
    p3.font.bold = True
    p3.font.color.rgb = PURE_WHITE
    p3.font.name = "Segoe UI"
    p4 = tf.add_paragraph()
    p4.text = "  - Simple LLM syntax errors (missing brackets, commas) cause validator parsing fails, leading to 15s timeout penalties."
    p4.font.size = Pt(13)
    p4.font.color.rgb = CHARCOAL_MUTED
    p4.font.name = "Segoe UI"

    p5 = tf.add_paragraph()
    p5.text = "• Code Execution Vulnerabilities & Key Theft Risks"
    p5.font.size = Pt(15)
    p5.font.bold = True
    p5.font.color.rgb = PURE_WHITE
    p5.font.name = "Segoe UI"
    p6 = tf.add_paragraph()
    p6.text = "  - Running unverified validator task scripts directly exposes node API keys and hotkey wallet mnemonics to RCE hacks."
    p6.font.size = Pt(13)
    p6.font.color.rgb = CHARCOAL_MUTED
    p6.font.name = "Segoe UI"

    p7 = tf.add_paragraph()
    p7.text = "• Boilerplate Token Bloat driving Latency and Cost"
    p7.font.size = Pt(15)
    p7.font.bold = True
    p7.font.color.rgb = PURE_WHITE
    p7.font.name = "Segoe UI"
    p8 = tf.add_paragraph()
    p8.text = "  - Chatty conversational filler output from LLMs increases payload sizes, bloating bandwidth and incurring latency penalties."
    p8.font.size = Pt(13)
    p8.font.color.rgb = CHARCOAL_MUTED
    p8.font.name = "Segoe UI"

    # ----------------------------------------------------
    # Slide 3: The Solution (Technical Breakthrough)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(content_layout)
    set_premium_background(slide)
    
    title_box = slide.shapes.title
    title_box.text = "02. The Solution: Hardware-Level System Optimization"
    title_box.text_frame.paragraphs[0].font.size = Pt(28)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = CYAN_ELECTRIC
    title_box.text_frame.paragraphs[0].font.name = "Segoe UI"
    
    tf = slide.placeholders[1].text_frame
    tf.clear()
    
    p1 = tf.add_paragraph()
    p1.text = "• Rust-Based Thread-Safe VRAM Scheduler"
    p1.font.size = Pt(15)
    p1.font.bold = True
    p1.font.color.rgb = PURE_WHITE
    p1.font.name = "Segoe UI"
    p2 = tf.add_paragraph()
    p2.text = "  - Actively throttles and queues over-limit requests, keeping VRAM stable and reducing OOM crash rate to 0%."
    p2.font.size = Pt(13)
    p2.font.color.rgb = CHARCOAL_MUTED
    p2.font.name = "Segoe UI"

    p3 = tf.add_paragraph()
    p3.text = "• Dynamic Schema SapParser"
    p3.font.size = Pt(15)
    p3.font.bold = True
    p3.font.color.rgb = PURE_WHITE
    p3.font.name = "Segoe UI"
    p4 = tf.add_paragraph()
    p4.text = "  - Automatically recovers and maps JSON structures to expected synapse schemas on-the-fly under microseconds."
    p4.font.size = Pt(13)
    p4.font.color.rgb = CHARCOAL_MUTED
    p4.font.name = "Segoe UI"

    p5 = tf.add_paragraph()
    p5.text = "• Deterministic AST Code Sanitizer & Mapper"
    p5.font.size = Pt(15)
    p5.font.bold = True
    p5.font.color.rgb = PURE_WHITE
    p5.font.name = "Segoe UI"
    p6 = tf.add_paragraph()
    p6.text = "  - Scans code tree nodes locally to intercept forbidden system imports and exec/eval operations before execution under 1µs."
    p6.font.size = Pt(13)
    p6.font.color.rgb = CHARCOAL_MUTED
    p6.font.name = "Segoe UI"

    p7 = tf.add_paragraph()
    p7.text = "• Boilerplate Token Compressor"
    p7.font.size = Pt(15)
    p7.font.bold = True
    p7.font.color.rgb = PURE_WHITE
    p7.font.name = "Segoe UI"
    p8 = tf.add_paragraph()
    p8.text = "  - Strips out conversational noise and filler phrases on-device, compressing payload sizes by 65-75% for speed."
    p8.font.size = Pt(13)
    p8.font.color.rgb = CHARCOAL_MUTED
    p8.font.name = "Segoe UI"

    # ----------------------------------------------------
    # Slide 4: Before vs After (Sleek Comparison Table)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_premium_background(slide)
    
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.8))
    tf = tx_box.text_frame
    p = tf.paragraphs[0]
    p.text = "03. Before vs After: The Paradigm Shift"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = CYAN_ELECTRIC
    p.font.name = "Segoe UI"
    
    rows = 7
    cols = 4
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.33)
    height = Inches(5.0)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    table.columns[0].width = Inches(2.3)  # Metric
    table.columns[1].width = Inches(3.0)  # Before
    table.columns[2].width = Inches(3.5)  # After
    table.columns[3].width = Inches(3.5)  # Impact
    
    headers = ["Metric / Benchmark", "Before (Original Python)", "After (Surfclaw 2.0 Rust OS)", "Business Impact (Value Prop)"]
    data = [
        ["Avg. Execution Latency", "1.76s - 2.10s (GIL Bottleneck)", "0.42s - 0.60s (GIL Bypassed)", "400%+ Faster Execution Speed"],
        ["VRAM Resource Control", "Unmanaged peaks, causing OOM crashes", "Safe queuing, resource budget tracking", "0% Node Downtime Achieved"],
        ["Security Boot Latency", "Docker/VM initialization: 2.50s - 5.00s", "Firecracker UDS socket boot: 0.12s (120ms)", "95%+ Security Overhead Reduction"],
        ["JSON Parse Resiliency", "Output schema errors prompt 15s timeout", "SapParser real-time JSON repair", "0% Format-related Penalty Rate"],
        ["Mining Rewards Yield", "Low weights due to execution delays", "Consistently high weights, maximizing emissions", "Maximized Delegated Staking APY (Optimum Yields)"],
        ["Operational Experience", "Sleepless nights monitoring node crashes", "Passive node maintenance (Sleep soundly)", "Zero Operational Stress & Fatigue"]
    ]
    
    # Format Headers
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = CYAN_ELECTRIC
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = SPACE_BLACK
            p.font.name = "Segoe UI"
            p.alignment = PP_ALIGN.CENTER
            
    # Format Body
    for row_idx, row_data in enumerate(data):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = text
            cell.fill.solid()
            # Alternate row colors
            cell.fill.fore_color.rgb = RGBColor(17, 20, 24) if row_idx % 2 == 0 else RGBColor(22, 27, 34)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = PURE_WHITE
                p.font.name = "Segoe UI"
                if col_idx == 3 or (col_idx == 2 and "0.42s" in text) or (col_idx == 2 and "0%" in text) or (col_idx == 2 and "120ms" in text):
                    p.font.bold = True
                    p.font.color.rgb = NEON_GREEN
                p.alignment = PP_ALIGN.LEFT if col_idx != 0 else PP_ALIGN.CENTER

    # ----------------------------------------------------
    # Slide 5: Core Strengths (Pros & Cons)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(content_layout)
    set_premium_background(slide)
    
    title_box = slide.shapes.title
    title_box.text = "04. Core Strengths & Technical Mitigation"
    title_box.text_frame.paragraphs[0].font.size = Pt(28)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = CYAN_ELECTRIC
    title_box.text_frame.paragraphs[0].font.name = "Segoe UI"
    
    tf = slide.placeholders[1].text_frame
    tf.clear()
    
    p1 = tf.add_paragraph()
    p1.text = "👍 Core Strengths (Pros)"
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = NEON_GREEN
    p1.font.name = "Segoe UI"
    
    p2 = tf.add_paragraph()
    p2.text = "  - High ROI Efficiency: Maximize RTX 4090 performance without investing in premium enterprise GPUs.\n" \
             "  - Viral Adoption Loop: Competitors are forced to adopt Surfclaw to protect their emission stakes.\n" \
             "  - Zero Maintenance Fatigue: 24/7 node stability ends midnight server restart routines."
    p2.font.size = Pt(13)
    p2.font.color.rgb = PURE_WHITE
    p2.font.name = "Segoe UI"
    
    p3 = tf.add_paragraph()
    p3.text = "👎 Technical Considerations & Mitigation (Cons & Mitigation)"
    p3.font.size = Pt(16)
    p3.font.bold = True
    p3.font.color.rgb = WARNING_RED
    p3.font.name = "Segoe UI"
    
    p4 = tf.add_paragraph()
    p4.text = "  - Virtualization Hardware Required: Host CPU must support Linux KVM (Fully met on all modern GPU miners).\n" \
             "  - Initial Rust Setup: Requires compilation (Bypassed with our one-click automated setup.bat / setup.sh scripts)."
    p4.font.size = Pt(13)
    p4.font.color.rgb = PURE_WHITE
    p4.font.name = "Segoe UI"

    # ----------------------------------------------------
    # Slide 6: Market Expansion & Universal DePIN Compatibility
    # ----------------------------------------------------
    slide = prs.slides.add_slide(content_layout)
    set_premium_background(slide)
    
    title_box = slide.shapes.title
    title_box.text = "05. Market Expansion: Universal DePIN Compatibility"
    title_box.text_frame.paragraphs[0].font.size = Pt(28)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = CYAN_ELECTRIC
    title_box.text_frame.paragraphs[0].font.name = "Segoe UI"
    
    tf = slide.placeholders[1].text_frame
    tf.clear()
    
    p1 = tf.add_paragraph()
    p1.text = "• Akash Network & io.net (Decentralized GPU Clouds)"
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = PURE_WHITE
    p1.font.name = "Segoe UI"
    p2 = tf.add_paragraph()
    p2.text = "  - Application: Prevents VRAM OOM container crashes under concurrent load spikes, securing 99.9% uptime and eliminating cost leaks on rented server instances."
    p2.font.size = Pt(13)
    p2.font.color.rgb = CHARCOAL_MUTED
    p2.font.name = "Segoe UI"

    p3 = tf.add_paragraph()
    p3.text = "• Render Network & Clore.ai (Node Provider Protection)"
    p3.font.size = Pt(16)
    p3.font.bold = True
    p3.font.color.rgb = PURE_WHITE
    p3.font.name = "Segoe UI"
    p4 = tf.add_paragraph()
    p4.text = "  - Application: Hardware-level microVM encapsulation completely isolates untrusted model execution, protecting host devices from privilege escalations and wallet thefts."
    p4.font.size = Pt(13)
    p4.font.color.rgb = CHARCOAL_MUTED
    p4.font.name = "Segoe UI"

    p5 = tf.add_paragraph()
    p5.text = "• TAM (Total Addressable Market) Expansion"
    p5.font.size = Pt(16)
    p5.font.bold = True
    p5.font.color.rgb = PURE_WHITE
    p5.font.name = "Segoe UI"
    p6 = tf.add_paragraph()
    p6.text = "  - Evolution: Scale from a Bittensor acceleration client to the default secure middleware OS for the entire $15B decentralized GPU cluster market."
    p6.font.size = Pt(13)
    p6.font.color.rgb = CHARCOAL_MUTED
    p6.font.name = "Segoe UI"

    # ----------------------------------------------------
    # Slide 7: Legal Disclaimer
    # ----------------------------------------------------
    slide = prs.slides.add_slide(content_layout)
    set_premium_background(slide)
    
    title_box = slide.shapes.title
    title_box.text = "Legal Disclaimer"
    title_box.text_frame.paragraphs[0].font.size = Pt(28)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = CHARCOAL_MUTED
    title_box.text_frame.paragraphs[0].font.name = "Segoe UI"
    
    tf = slide.placeholders[1].text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "The development team provides software on an open-source basis only and offers no financial services.\n\n" \
             "The projected Staking APY mentioned in this document is an estimate based on the decentralized consensus rules and current dynamics of the Bittensor network. It does not represent a guarantee of interest rates or principal protection.\n\n" \
             "All staking transactions and reward distributions are processed autonomously by the on-chain smart contracts of the decentralized network. The project entity does not hold custody of user assets or distribute yields directly."
    p.font.size = Pt(13)
    p.font.color.rgb = CHARCOAL_MUTED
    p.font.name = "Segoe UI"

    # Save presentation
    filename = "Surfclaw_Pitch_Deck.pptx"
    prs.save(filename)
    print(f"[SUCCESS] PowerPoint Pitch Deck generated successfully: {os.path.abspath(filename)}")

if __name__ == "__main__":
    create_pitch_deck()
